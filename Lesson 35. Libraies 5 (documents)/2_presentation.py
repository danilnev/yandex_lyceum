from pptx import Presentation

doc = Presentation()
layout = doc.slide_layouts[2]
data = {
    'randint': '''Help on method randint in module random: 
    randint(a, b) method of random.Random instance
    Return random integer in range [a, b], including both end points.''',
    'randrange': '''Help on method randrange in module random:

randrange(start, stop=None, step=1) method of random.Random instance
    Choose a random item from range(stop) or range(start, stop[, step]).
    
    Roughly equivalent to ``choice(range(start, stop, step))`` but
    supports arbitrarily large ranges and is optimized for common cases.
    ''',
    'randbytes': '''Help on method randbytes in module random:

randbytes(n) method of random.Random instance
    Generate n random bytes.''',
    'choice': '''Help on method choice in module random:

choice(seq) method of random.Random instance
    Choose a random element from a non-empty sequence.''',
    'choices': '''Help on method choices in module random:

choices(population, weights=None, *, cum_weights=None, k=1) method of random.Random instance
    Return a k sized list of population elements chosen with replacement.
    
    If the relative weights or cumulative weights are not specified,
    the selections are made with equal probability.'''
}
for key in data:
    slide = doc.slides.add_slide(layout)
    slide.placeholders[0].text = key

    slide.placeholders[1].text = data[key]
    ph = slide.placeholders
    for p in ph:
        p.text_frame.fit_text(font_family='Courier New')
doc.save('doc.pptx')
